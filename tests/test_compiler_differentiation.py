"""Lock in that the three compiler presets produce visibly different decks.

If two presets produce structurally identical output, the compiler architecture
isn't earning its keep — these tests fail and force you to fix it.
"""
from collections import namedtuple

from pptx import Presentation

from tests.conftest import make_card

Shape = namedtuple("Shape", "n_slides n_pictures n_textboxes has_notes")


def _shape(path) -> Shape:
    prs = Presentation(str(path))
    pics = sum(1 for s in prs.slides for shp in s.shapes if shp.shape_type == 13)
    tx = sum(1 for s in prs.slides for shp in s.shapes if shp.has_text_frame)
    notes_present = any(
        s.notes_slide.notes_text_frame.text.strip()
        for s in prs.slides
    )
    return Shape(len(prs.slides), pics, tx, notes_present)


def _seed_with_corpus(mock_scryfall, mock_image_downloader):
    """Seed enough cards that all preset features get exercised."""
    from mtgslider import themes
    themes.create_theme(
        "Books", short_description="Cards depicting books, scrolls, and writing.", aliases=["scrolls", "tomes"],
    )
    cards = [
        make_card("Library of Alexandria", id="id-loa", mana_cost=""),
        make_card("Jace's Archivist", id="id-ja", mana_cost="{1}{U}{U}"),
        make_card("Whispering Madness", id="id-wm", mana_cost="{2}{U}{B}"),
        make_card("Ancestral Vision", id="id-av", mana_cost=""),
        make_card("Time Spiral", id="id-ts", mana_cost="{4}{U}{U}"),
        make_card("Snapcaster Mage", id="id-sm", mana_cost="{1}{U}"),
        make_card("Brainstorm", id="id-bs", mana_cost="{U}"),
        make_card("Counterspell", id="id-cs", mana_cost="{U}{U}"),
        make_card("Force of Will", id="id-fow", mana_cost="{3}{U}{U}"),
        make_card("Ponder", id="id-p", mana_cost="{U}"),
    ]
    mock_scryfall.searches["o:book"] = cards
    themes.find_cards_by_query("books", "o:book")
    themes.mark_card("books", "Library of Alexandria", include_status="exemplar")
    themes.mark_card("books", "Jace's Archivist", include_status="include")
    themes.mark_card("books", "Whispering Madness", include_status="include")
    themes.mark_card("books", "Ancestral Vision", include_status="include")
    themes.mark_card("books", "Time Spiral", include_status="include")
    themes.mark_card("books", "Snapcaster Mage", include_status="include")
    themes.mark_card("books", "Brainstorm", include_status="include")
    themes.fetch_images_for_theme("books", variant="normal")


def test_three_presets_produce_distinct_shapes(mock_scryfall, mock_image_downloader):
    _seed_with_corpus(mock_scryfall, mock_image_downloader)
    from mtgslider.slideshow.compiler.generator import build

    shapes = {s: _shape(build("books", style=s)) for s in ("rhystic", "sierkovitz", "teaching")}
    # Reduce shapes to (n_slides, n_pictures) — those are the load-bearing differences a viewer notices.
    fingerprints = {k: (v.n_slides, v.n_pictures) for k, v in shapes.items()}
    assert len(set(fingerprints.values())) == 3, (
        f"Expected three distinct (slides, pictures) fingerprints, got {fingerprints}. "
        "Two or more presets are producing structurally identical output."
    )


def test_sierkovitz_includes_data_slide_others_do_not(mock_scryfall, mock_image_downloader):
    _seed_with_corpus(mock_scryfall, mock_image_downloader)
    from mtgslider.slideshow.compiler.generator import build
    from mtgslider.slideshow.compiler.planner import plan
    from mtgslider.slideshow.compiler.presets import resolve
    from mtgslider.packet import build_packet

    packet = build_packet("books")
    sier_units = plan(packet, resolve("sierkovitz"))
    rhys_units = plan(packet, resolve("rhystic"))
    teach_units = plan(packet, resolve("teaching"))

    sier_types = [u.slide_type for u in sier_units]
    rhys_types = [u.slide_type for u in rhys_units]
    teach_types = [u.slide_type for u in teach_units]
    assert "data" in sier_types
    assert "data" not in rhys_types
    assert "data" not in teach_types


def test_rhystic_includes_visual_motif_sierkovitz_does_not(mock_scryfall, mock_image_downloader):
    _seed_with_corpus(mock_scryfall, mock_image_downloader)
    from mtgslider.slideshow.compiler.planner import plan
    from mtgslider.slideshow.compiler.presets import resolve
    from mtgslider.packet import build_packet

    packet = build_packet("books")
    rhys = [u.slide_type for u in plan(packet, resolve("rhystic"))]
    sier = [u.slide_type for u in plan(packet, resolve("sierkovitz"))]
    assert "visual_motif" in rhys
    assert "visual_motif" not in sier


def test_teaching_includes_takeaways(mock_scryfall, mock_image_downloader):
    _seed_with_corpus(mock_scryfall, mock_image_downloader)
    from mtgslider.slideshow.compiler.planner import plan
    from mtgslider.slideshow.compiler.presets import resolve
    from mtgslider.packet import build_packet

    packet = build_packet("books")
    teach = [u.slide_type for u in plan(packet, resolve("teaching"))]
    rhys = [u.slide_type for u in plan(packet, resolve("rhystic"))]
    assert "key_takeaways" in teach
    assert "key_takeaways" not in rhys


def test_all_presets_emit_speaker_notes(mock_scryfall, mock_image_downloader):
    _seed_with_corpus(mock_scryfall, mock_image_downloader)
    from mtgslider.slideshow.compiler.generator import build

    for s in ("rhystic", "sierkovitz", "teaching"):
        sh = _shape(build("books", style=s))
        assert sh.has_notes, f"preset {s!r} produced a deck with no speaker notes"


def test_speaker_notes_voice_differs_across_presets(mock_scryfall, mock_image_downloader):
    """The notes for the same slide_type should not be character-identical across voices."""
    _seed_with_corpus(mock_scryfall, mock_image_downloader)
    from mtgslider.slideshow.compiler.notes import for_unit

    theme = {"canonical_name": "Books", "aliases": [], "short_description": "books theme"}
    a = for_unit({"card": {"card_name": "Brainstorm"}, "claim": "iconic blue draw"}, "card_spotlight", mode="rich", theme=theme)
    b = for_unit({"card": {"card_name": "Brainstorm"}, "claim": "iconic blue draw"}, "card_spotlight", mode="terse", theme=theme)
    c = for_unit({"card": {"card_name": "Brainstorm"}, "claim": "iconic blue draw"}, "card_spotlight", mode="heuristic", theme=theme)
    assert a and b and c
    assert len({a, b, c}) == 3, "expected three distinct voices for the same card spotlight"


def test_evidence_mode_annotated_adds_footer_text(mock_scryfall, mock_image_downloader):
    """In annotated mode, card spotlight slides carry a 'via Scryfall' footer; clean mode does not."""
    _seed_with_corpus(mock_scryfall, mock_image_downloader)
    from mtgslider.slideshow.compiler.generator import build

    def all_text(path):
        prs = Presentation(str(path))
        out = []
        for s in prs.slides:
            for shp in s.shapes:
                if shp.has_text_frame:
                    out.append(shp.text_frame.text)
        return "\n".join(out)

    rhys = all_text(build("books", style="rhystic"))      # clean — no footer
    sier = all_text(build("books", style="sierkovitz"))   # annotated
    assert "via Scryfall" not in rhys
    assert "via Scryfall" in sier


def test_v1_does_not_emit_speaker_notes(mock_scryfall, mock_image_downloader):
    """v1_template stays a deterministic template; speaker notes are a compiler-only feature."""
    _seed_with_corpus(mock_scryfall, mock_image_downloader)
    from mtgslider.slideshow.v1_template.generator import build

    sh = _shape(build("books"))
    assert not sh.has_notes


def test_compiler_distinct_from_v1(mock_scryfall, mock_image_downloader):
    """The compiler should not produce a deck structurally identical to v1 in any preset."""
    _seed_with_corpus(mock_scryfall, mock_image_downloader)
    from mtgslider.slideshow.compiler.generator import build as build_compiler
    from mtgslider.slideshow.v1_template.generator import build as build_v1

    v1_shape = _shape(build_v1("books"))
    for s in ("rhystic", "sierkovitz", "teaching"):
        compiler_shape = _shape(build_compiler("books", style=s))
        assert (
            (v1_shape.n_slides, v1_shape.n_pictures, v1_shape.n_textboxes)
            != (compiler_shape.n_slides, compiler_shape.n_pictures, compiler_shape.n_textboxes)
        ), f"compiler preset {s!r} produced a deck structurally identical to v1"
