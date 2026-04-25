"""Render ArgumentUnits to .pptx via the layout registry.

Preset-aware: the same ArgumentUnit can render to different layouts depending on
the active preset. This is where the presets earn their keep.
"""
from __future__ import annotations

from datetime import datetime
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .layouts import Layout, LAYOUT_REGISTRY
from .slide_types import ArgumentUnit

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

_GREY = RGBColor(0x55, 0x55, 0x55)
_MUTED = RGBColor(0x99, 0x99, 0x99)
_RED = RGBColor(0xB0, 0x00, 0x00)


def render(units: list[ArgumentUnit], preset: dict) -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for unit in units:
        layout = _pick_layout(unit, preset)
        _render_unit(prs, unit, layout, preset)
    return prs


def _pick_layout(unit: ArgumentUnit, preset: dict) -> Layout:
    # Preset-specific overrides first.
    if unit.slide_type == "cluster" and preset.get("use_dense_clusters"):
        return LAYOUT_REGISTRY["cluster_grid_dense"]
    if unit.slide_type == "card_spotlight" and preset.get("density") == "low":
        return LAYOUT_REGISTRY["single_image_focus"]
    if unit.slide_type == "card_spotlight" and preset.get("density") == "high":
        return LAYOUT_REGISTRY["split_image_text"]
    # Generic fallback: first compatible layout.
    candidates = [l for l in LAYOUT_REGISTRY.values() if unit.slide_type in l.compatible_slide_types]
    return candidates[0] if candidates else LAYOUT_REGISTRY["text_only"]


def _render_unit(prs: Presentation, unit: ArgumentUnit, layout: Layout, preset: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    handler = _RENDERERS.get(layout.name, _render_text_only)
    handler(slide, unit, layout, preset)
    _apply_evidence_footer(slide, unit, preset)
    _apply_speaker_notes(slide, unit)


# ---------- per-layout handlers ----------

def _render_text_only(slide, unit, layout, preset):
    title = unit.content.get("title", "")
    body_parts: list[str] = []
    for k in ("subtitle", "definition", "claim"):
        v = unit.content.get(k)
        if v:
            body_parts.append(str(v))
    examples = unit.content.get("examples")
    if examples:
        body_parts.append("Examples: " + ", ".join(str(e) for e in examples))
    sources = unit.content.get("sources")
    if sources:
        body_parts.append("\n".join(f"• {s}" for s in sources))

    _add_text(slide, title, left=Inches(0.7), top=Inches(0.7), width=Inches(12), height=Inches(1.2),
              size=int(36 * layout.typography_scale), bold=True)
    if body_parts:
        body = "\n\n".join(body_parts)
        body = _trim_words(body, preset.get("max_words_per_slide", 60))
        _add_text(slide, body,
                  left=Inches(0.7), top=Inches(2.2), width=Inches(12), height=Inches(4.8),
                  size=int(20 * layout.typography_scale))


def _render_split_image_text(slide, unit, layout, preset):
    _render_card_body(slide, unit, preset, text_left_in=5.5, text_width_in=7.3)


def _render_single_image_focus(slide, unit, layout, preset):
    # Low-density variant: bigger image (wider horizontal), tighter caption on the right.
    _render_card_body(slide, unit, preset, text_left_in=6.2, text_width_in=6.6, image_height_in=6.4)


def _render_card_body(slide, unit, preset, *, text_left_in: float, text_width_in: float, image_height_in: float = 6.0):
    card = unit.content.get("card") or {}
    img = unit.content.get("image_path") or card.get("image_path")
    name = card.get("card_name", unit.content.get("title", ""))
    type_line = card.get("type_line", "")
    mana = card.get("mana_cost", "")
    label = unit.content.get("label") or ""
    claim = unit.content.get("claim") or card.get("link_reason") or ""
    note = card.get("note") or ""
    if img:
        try:
            slide.shapes.add_picture(img, Inches(0.5), Inches(0.7), height=Inches(image_height_in))
        except (FileNotFoundError, OSError):
            pass
    _add_text(slide, name,
              left=Inches(text_left_in), top=Inches(0.7), width=Inches(text_width_in), height=Inches(0.8),
              size=32, bold=True)
    _add_text(slide, f"{type_line}   {mana}",
              left=Inches(text_left_in), top=Inches(1.6), width=Inches(text_width_in), height=Inches(0.5),
              size=16, color=_GREY)
    if label:
        _add_text(slide, label,
                  left=Inches(text_left_in), top=Inches(2.2), width=Inches(3), height=Inches(0.4),
                  size=12, bold=True, color=_RED)
    body = "\n\n".join(s for s in (claim, note) if s)
    if body:
        body = _trim_words(body, preset.get("max_words_per_slide", 60))
        _add_text(slide, body,
                  left=Inches(text_left_in), top=Inches(2.8), width=Inches(text_width_in), height=Inches(4),
                  size=16)


def _render_grid_2x2(slide, unit, layout, preset):
    _render_grid(slide, unit, cols=4, rows=2, cell_h_in=2.6, image_h_in=2.1)


def _render_cluster_grid_dense(slide, unit, layout, preset):
    _render_grid(slide, unit, cols=3, rows=3, cell_h_in=1.9, image_h_in=1.5)


def _render_grid(slide, unit, *, cols: int, rows: int, cell_h_in: float, image_h_in: float):
    title = unit.content.get("title", "")
    cards = unit.content.get("cards") or []
    _add_text(slide, title,
              left=Inches(0.5), top=Inches(0.3), width=Inches(12), height=Inches(0.6),
              size=26, bold=True)
    cell_w = (SLIDE_W - Inches(1)) / cols
    cell_h = Inches(cell_h_in)
    for i, card in enumerate(cards[: cols * rows]):
        r, c = divmod(i, cols)
        left = Inches(0.5) + cell_w * c
        top = Inches(1.1) + cell_h * r
        if card.get("image_path"):
            try:
                slide.shapes.add_picture(card["image_path"], left, top, height=Inches(image_h_in))
            except (FileNotFoundError, OSError):
                pass
        _add_text(slide, card.get("card_name", ""),
                  left=left, top=top + Inches(image_h_in), width=cell_w, height=Inches(0.35),
                  size=11, bold=True)


def _render_data_chart_text(slide, unit, layout, preset):
    title = unit.content.get("title", "")
    stats = unit.content.get("stats") or {}
    claim = unit.content.get("claim") or ""

    _add_text(slide, title,
              left=Inches(0.5), top=Inches(0.4), width=Inches(12), height=Inches(0.8),
              size=30, bold=True)

    # Stats table on left as text rows.
    lines: list[str] = [f"Cards: {stats.get('card_count', 0)}"]
    by_color = stats.get("by_color") or {}
    if by_color:
        lines.append("")
        lines.append("By color:")
        for k, v in by_color.items():
            lines.append(f"  {k}: {v}")
    by_type = stats.get("by_type") or {}
    if by_type:
        lines.append("")
        lines.append("By type:")
        for k, v in by_type.items():
            lines.append(f"  {k}: {v}")
    by_mv = stats.get("by_mana_value") or {}
    if by_mv:
        lines.append("")
        lines.append("By mana value:")
        for k, v in sorted(by_mv.items()):
            lines.append(f"  {k}: {v}")
    _add_text(slide, "\n".join(lines),
              left=Inches(0.6), top=Inches(1.3), width=Inches(5.5), height=Inches(5.5),
              size=14)

    if claim:
        _add_text(slide, claim,
                  left=Inches(6.8), top=Inches(1.3), width=Inches(6), height=Inches(5),
                  size=18)


def _render_takeaways(slide, unit, layout, preset):
    title = unit.content.get("title", "")
    takeaways = unit.content.get("takeaways") or []
    _add_text(slide, title,
              left=Inches(0.7), top=Inches(0.6), width=Inches(12), height=Inches(0.9),
              size=32, bold=True)
    numbered = "\n\n".join(f"{i + 1}. {t}" for i, t in enumerate(takeaways))
    _add_text(slide, numbered,
              left=Inches(0.7), top=Inches(1.8), width=Inches(12), height=Inches(5),
              size=int(22 * layout.typography_scale))


_RENDERERS = {
    "text_only": _render_text_only,
    "split_image_text": _render_split_image_text,
    "single_image_focus": _render_single_image_focus,
    "grid_2x2": _render_grid_2x2,
    "cluster_grid_dense": _render_cluster_grid_dense,
    "data_chart_text": _render_data_chart_text,
    "takeaways": _render_takeaways,
}


# ---------- speaker notes + evidence footer ----------

def _apply_speaker_notes(slide, unit: ArgumentUnit) -> None:
    if not unit.speaker_notes:
        return
    nslide = slide.notes_slide
    tf = nslide.notes_text_frame
    tf.text = unit.speaker_notes


def _apply_evidence_footer(slide, unit: ArgumentUnit, preset: dict) -> None:
    if preset.get("evidence_mode") != "annotated":
        return
    if unit.slide_type == "card_spotlight":
        card = unit.content.get("card") or {}
        set_code = (card.get("image_source_url") or "")
        # Prefer Scryfall provenance already on the card entry.
        pieces = ["via Scryfall"]
        if card.get("scryfall_id"):
            pieces.append(card["scryfall_id"][:8])
        if card.get("image_variant"):
            pieces.append(f"{card['image_variant']}")
        footer = " · ".join(pieces)
    elif unit.slide_type in ("cluster", "visual_motif"):
        cards = unit.content.get("cards") or []
        footer = f"via Scryfall · {len(cards)} cards"
    elif unit.slide_type == "data":
        footer = "via Scryfall · aggregated from curated set"
    else:
        return
    _add_text(slide, footer,
              left=Inches(0.4), top=Inches(7.05), width=Inches(12.5), height=Inches(0.3),
              size=9, color=_MUTED)


# ---------- helpers ----------

def _add_text(slide, text, *, left, top, width, height, size=18, bold=False, align=PP_ALIGN.LEFT, color=None):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _trim_words(text: str, max_words: int) -> str:
    words = text.split()
    if len(words) <= max_words:
        return text
    return " ".join(words[:max_words]) + " …"
